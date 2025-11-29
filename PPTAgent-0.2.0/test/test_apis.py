from bs4 import BeautifulSoup
from pptx import Presentation

from pptagent.apis import (
    API_TYPES,
    CodeExecutor,
    markdown,
    process_element,
    replace_para,
)
from test.conftest import test_config


def test_api_docs():
    executor = CodeExecutor(3)
    docs = executor.get_apis_docs(API_TYPES.Agent.value)
    assert len(docs) > 0


def test_replace_para():
    text = "这是一个**加粗和*斜体*文本**，还有*斜体和`Code def a+b`*，~~删除~~，[链接](http://example.com)"
    prs = Presentation(test_config.ppt)
    slide = prs.slides[0]
    replace_para(0, text, slide.shapes[0])
    runs = slide.shapes[0].text_frame.paragraphs[0].runs
    assert runs[1].font.bold
    assert runs[2].font.bold and runs[2].font.italic
    assert runs[6].font.name == "Consolas"
    assert runs[8].font.strikethrough
    assert runs[10].hyperlink.address == "http://example.com"




'''
测你的系统是否把 Markdown 列表
转为普通文本段落
（而非 HTML 列表）
'''
def test_list_parsing():
    text = """
    - 项目1
    - 项目2

    1. 项目1
    2. 项目2
    """
    
    '''
        markdown = create_markdown(renderer=SlideRenderer(), plugins=["strikethrough"])
        from mistune import HTMLRenderer, create_markdown
        markdown可以理解为操作富文本的工具，它可以将富文本转化为html，然后再用专业的工具读取转化为pptx能够接受的_Run对象
    '''
    html = markdown(text).strip()
    soup = BeautifulSoup(html, "html.parser")
    blocks = process_element(soup)
    assert len(blocks) == 1
    assert "ol" not in html and "ul" not in html
